from __future__ import annotations

import json
import re
from pathlib import Path
from typing import Any

from agent.rag.config import RagConfig
from agent.rag.schemas import Citation, EvidenceItem, RagPackage


SUPPORTED_EXTENSIONS = {".md", ".txt", ".docx", ".xlsx", ".pdf", ".pptx"}
CHUNK_CHARS = 900
CHUNK_OVERLAP = 120


class SimpleResourceRetriever:
    """Lightweight fallback retriever for local resource files.

    It first tries real docstore JSON when available. If the copied index files
    are Git LFS pointers, it searches the source files directly.
    """

    def __init__(self, config: RagConfig | None = None) -> None:
        self.config = config or RagConfig.from_env()

    def retrieve(self, queries: list[str]) -> RagPackage:
        clean_queries = [item.strip() for item in queries if item and item.strip()]
        query = "；".join(clean_queries)
        warnings: list[str] = []

        chunks = self._load_index_chunks(warnings)
        if not chunks:
            chunks = self._load_source_chunks(clean_queries, warnings)

        scored: list[EvidenceItem] = []
        for item in chunks:
            score = _score_text(clean_queries, item.text, item.source_file)
            if score >= self.config.min_score:
                scored.append(item.model_copy(update={"score": round(score, 4)}))

        scored.sort(key=lambda item: item.score, reverse=True)
        evidence = _select_diverse_evidence(scored, self.config.top_k)
        citations = [_citation_for(item) for item in evidence]
        confidence = round(max((item.score for item in evidence), default=0.0), 4)
        next_action = "use_as_grounded_context" if evidence else "need_more_evidence"

        if not evidence:
            warnings.append("evidence_not_found")

        return RagPackage(
            query=query,
            answer=_extractive_answer(evidence),
            evidence=evidence,
            citations=citations,
            confidence=confidence,
            warnings=list(dict.fromkeys(warnings)),
            next_action=next_action,
        )

    def _load_index_chunks(self, warnings: list[str]) -> list[EvidenceItem]:
        docstore_path = self.config.index_dir / "docstore.json"
        if not docstore_path.exists():
            warnings.append(f"index_docstore_missing:{docstore_path}")
            return []

        text = docstore_path.read_text(encoding="utf-8", errors="ignore")
        if text.startswith("version https://git-lfs.github.com/spec"):
            warnings.append("index_docstore_is_lfs_pointer")
            return []

        try:
            value = json.loads(text)
        except json.JSONDecodeError:
            warnings.append("index_docstore_json_invalid")
            return []

        chunks = []
        for index, raw_node in enumerate(_iter_docstore_nodes(value), start=1):
            node = _node_payload(raw_node)
            node_text = _node_text(node)
            if not node_text:
                continue
            metadata = node.get("metadata") if isinstance(node, dict) else {}
            metadata = metadata if isinstance(metadata, dict) else {}
            source_file = str(
                metadata.get("source_file")
                or metadata.get("file_name")
                or metadata.get("source_doc")
                or "docstore"
            )
            chunks.append(
                EvidenceItem(
                    source_file=source_file,
                    chunk_id=str(node.get("id_") or node.get("node_id") or f"docstore_{index}"),
                    text=node_text,
                    file_type=str(metadata.get("file_type") or ""),
                    page_label=_optional_str(metadata.get("page_label") or metadata.get("page")),
                    metadata=metadata,
                )
            )
        return chunks

    def _load_source_chunks(self, queries: list[str], warnings: list[str]) -> list[EvidenceItem]:
        files: list[Path] = []
        seen_files: set[Path] = set()
        for source_dir in _source_dirs(self.config):
            if not source_dir.exists():
                warnings.append(f"source_dir_missing:{source_dir}")
                continue
            for item in source_dir.rglob("*"):
                if not item.is_file() or item.suffix.lower() not in SUPPORTED_EXTENSIONS:
                    continue
                resolved = item.resolve()
                if resolved in seen_files:
                    continue
                seen_files.add(resolved)
                files.append(item)
        files.sort(key=lambda path: _score_filename(queries, path.name), reverse=True)

        chunks: list[EvidenceItem] = []
        for path in files:
            try:
                sections = _read_source_sections(path, max_pdf_pages=self.config.max_pdf_pages)
            except Exception as exc:
                warnings.append(f"source_read_failed:{path.name}:{type(exc).__name__}")
                continue
            chunks_before = len(chunks)
            for section_index, section in enumerate(sections, start=1):
                chunks.extend(
                    _split_text(
                        section["text"],
                        source_file=path.name,
                        file_type=path.suffix.lower().lstrip("."),
                        page_label=_optional_str(section.get("page_label")),
                        metadata={
                            "source_path": str(path.resolve()),
                            "section_index": section_index,
                            **dict(section.get("metadata") or {}),
                        },
                    )
                )
            if len(chunks) == chunks_before:
                warnings.append(f"source_text_empty:{path.name}")
        return chunks


def _source_dirs(config: RagConfig) -> list[Path]:
    dirs: list[Path] = []
    seen: set[Path] = set()
    for source_dir in [config.source_dir, *getattr(config, "additional_source_dirs", ())]:
        key = source_dir.resolve() if source_dir.exists() else source_dir
        if key in seen:
            continue
        seen.add(key)
        dirs.append(source_dir)
    return dirs


def _iter_docstore_nodes(value: Any):
    if isinstance(value, dict):
        for key in ("docstore/data", "data", "docs", "nodes"):
            nested = value.get(key)
            if isinstance(nested, dict):
                yield from nested.values()
            elif isinstance(nested, list):
                yield from nested
        for item in value.values():
            if isinstance(item, dict) and ("text" in item or "metadata" in item):
                yield item


def _node_payload(node: Any) -> dict[str, Any]:
    if not isinstance(node, dict):
        return {}
    wrapped = node.get("__data__")
    if isinstance(wrapped, dict):
        return wrapped
    return node


def _node_text(node: Any) -> str:
    if not isinstance(node, dict):
        return ""
    for key in ("text", "content", "text_resource"):
        value = node.get(key)
        if isinstance(value, str):
            return _clean_text(value)
        if isinstance(value, dict) and isinstance(value.get("text"), str):
            return _clean_text(value["text"])
    return ""


def _read_source_text(path: Path, *, max_pdf_pages: int) -> str:
    """Compatibility helper returning all directly extracted source text."""
    return "\n".join(
        section["text"] for section in _read_source_sections(path, max_pdf_pages=max_pdf_pages)
    )


def _read_source_sections(path: Path, *, max_pdf_pages: int) -> list[dict[str, Any]]:
    """Read a source without embeddings while retaining useful source locators."""
    suffix = path.suffix.lower()
    if suffix in {".md", ".txt"}:
        return [{"text": path.read_text(encoding="utf-8", errors="ignore")}]
    if suffix == ".docx":
        from docx import Document

        document = Document(str(path))
        parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
        for table in document.tables:
            for row in table.rows:
                values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if values:
                    parts.append(" | ".join(values))
        return [{"text": "\n".join(parts)}]
    if suffix == ".xlsx":
        from openpyxl import load_workbook

        workbook = load_workbook(str(path), read_only=True, data_only=True)
        sections: list[dict[str, Any]] = []
        for sheet in workbook.worksheets:
            rows = []
            for row in sheet.iter_rows(values_only=True):
                values = [str(cell) for cell in row if cell is not None]
                if values:
                    rows.append(" | ".join(values))
            if rows:
                sections.append(
                    {
                        "text": "\n".join(rows),
                        "page_label": sheet.title,
                        "metadata": {"locator_type": "sheet"},
                    }
                )
        return sections
    if suffix == ".pdf":
        return _read_pdf_sections(path, max_pages=max_pdf_pages)
    if suffix == ".pptx":
        return _read_pptx_sections(path)
    return []


def _read_pdf_text(path: Path, *, max_pages: int) -> str:
    return "\n".join(section["text"] for section in _read_pdf_sections(path, max_pages=max_pages))


def _read_pdf_sections(path: Path, *, max_pages: int) -> list[dict[str, Any]]:
    try:
        import fitz

        with fitz.open(path) as document:
            page_count = document.page_count if max_pages <= 0 else min(document.page_count, max_pages)
            pages = [
                {"text": document[index].get_text("text"), "page_label": str(index + 1)}
                for index in range(page_count)
            ]
        return pages
    except ModuleNotFoundError:
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        selected_pages = reader.pages if max_pages <= 0 else reader.pages[:max_pages]
        return [
            {"text": page.extract_text() or "", "page_label": str(index)}
            for index, page in enumerate(selected_pages, start=1)
        ]


def _read_pptx_text(path: Path) -> str:
    return "\n".join(section["text"] for section in _read_pptx_sections(path))


def _read_pptx_sections(path: Path) -> list[dict[str, Any]]:
    from pptx import Presentation

    presentation = Presentation(str(path))
    sections: list[dict[str, Any]] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if values:
                        parts.append(" | ".join(values))
        if parts:
            sections.append(
                {
                    "text": "\n".join(parts),
                    "page_label": str(slide_index),
                    "metadata": {"locator_type": "slide"},
                }
            )
    return sections


def _split_text(
    text: str,
    *,
    source_file: str,
    file_type: str,
    page_label: str | None = None,
    metadata: dict[str, Any] | None = None,
) -> list[EvidenceItem]:
    cleaned = _clean_text(text)
    if not cleaned:
        return []
    chunks = []
    start = 0
    index = 1
    while start < len(cleaned):
        end = min(start + CHUNK_CHARS, len(cleaned))
        chunk = cleaned[start:end].strip()
        if chunk:
            chunks.append(
                EvidenceItem(
                    source_file=source_file,
                    chunk_id=(
                        f"{Path(source_file).stem}_"
                        f"{_safe_identifier(page_label or 'document')}_{index:04d}"
                    ),
                    text=chunk,
                    file_type=file_type,
                    page_label=page_label,
                    metadata=dict(metadata or {}),
                )
            )
        if end == len(cleaned):
            break
        start = max(end - CHUNK_OVERLAP, start + 1)
        index += 1
    return chunks


def _score_text(queries: list[str], text: str, source_file: str) -> float:
    haystack = _normalize_search_text(f"{source_file}\n{text}")
    weighted_terms = _weighted_query_terms(" ".join(queries))
    if not weighted_terms:
        return 0.0
    total_weight = sum(weighted_terms.values())
    matched_weight = sum(weight for term, weight in weighted_terms.items() if term in haystack)
    coverage = matched_weight / total_weight if total_weight else 0.0
    frequency = sum(min(haystack.count(term), 4) * weight for term, weight in weighted_terms.items())
    frequency_bonus = min(frequency / max(total_weight * 8, 1), 0.18)
    technical = [term for term in weighted_terms if re.search(r"[a-z]\d|\d[a-z]", term)]
    technical_bonus = min(sum(1 for term in technical if term in haystack) * 0.08, 0.24)
    filename = _normalize_search_text(source_file)
    filename_bonus = min(sum(weight for term, weight in weighted_terms.items() if term in filename) / 40, 0.12)
    return min(1.0, coverage * 0.68 + frequency_bonus + technical_bonus + filename_bonus)


def _score_filename(queries: list[str], filename: str) -> float:
    name = filename.lower()
    tokens = _query_tokens(" ".join(queries))
    return sum(1 for token in tokens if token.lower() in name)


def _query_tokens(text: str) -> list[str]:
    return list(_weighted_query_terms(text))


def _weighted_query_terms(text: str) -> dict[str, float]:
    normalized = _normalize_search_text(text)
    result: dict[str, float] = {}
    for token in re.findall(r"[a-z]+\d+[a-z0-9_]*|\d+[a-z]+[a-z0-9_]*|[a-z_]{2,}", normalized):
        result[token] = max(result.get(token, 0.0), 4.0 if any(char.isdigit() for char in token) else 2.0)
    for phrase in re.findall(r"[\u4e00-\u9fff]{2,}", normalized):
        if len(phrase) <= 8:
            result[phrase] = max(result.get(phrase, 0.0), 3.0)
        for size, weight in ((4, 2.0), (3, 1.5), (2, 1.0)):
            for index in range(0, max(len(phrase) - size + 1, 0)):
                term = phrase[index : index + size]
                result[term] = max(result.get(term, 0.0), weight)
    return result


def _search_chars(text: str) -> set[str]:
    return {char for char in text if "\u4e00" <= char <= "\u9fff"}


def _citation_for(item: EvidenceItem) -> Citation:
    label = item.source_file
    if item.page_label:
        locator_type = item.metadata.get("locator_type")
        if locator_type == "sheet":
            label = f"{label}，工作表 {item.page_label}"
        elif locator_type == "slide":
            label = f"{label}，第 {item.page_label} 张幻灯片"
        else:
            label = f"{label}，第 {item.page_label} 页"
    return Citation(source_file=item.source_file, chunk_id=item.chunk_id, label=label)


def _extractive_answer(evidence: list[EvidenceItem]) -> str:
    if not evidence:
        return ""
    return "\n".join(f"[{index}] {item.text[:220]}" for index, item in enumerate(evidence, start=1))


def _clean_text(text: str) -> str:
    return re.sub(r"\s+", " ", text).strip()


def _normalize_search_text(text: str) -> str:
    return text.lower().replace("ｇ", "g").replace("０", "0").replace("２", "2").replace("３", "3")


def _safe_identifier(value: str) -> str:
    return re.sub(r"[^0-9A-Za-z\u4e00-\u9fff_-]+", "_", value).strip("_") or "section"


def _select_diverse_evidence(scored: list[EvidenceItem], top_k: int) -> list[EvidenceItem]:
    if top_k <= 0:
        return []
    per_source_limit = max(2, (top_k + 1) // 2)
    selected: list[EvidenceItem] = []
    source_counts: dict[str, int] = {}
    for item in scored:
        if source_counts.get(item.source_file, 0) >= per_source_limit:
            continue
        selected.append(item)
        source_counts[item.source_file] = source_counts.get(item.source_file, 0) + 1
        if len(selected) >= top_k:
            break
    if len(selected) < top_k:
        selected_ids = {item.chunk_id for item in selected}
        selected.extend(item for item in scored if item.chunk_id not in selected_ids)
    return selected[:top_k]


def _optional_str(value: Any) -> str | None:
    return str(value) if value is not None else None
